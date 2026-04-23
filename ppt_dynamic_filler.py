#!/usr/bin/env python3
"""
按「列表嵌套字典」动态拼装 PPT：从模板中按占位符键集合匹配版式页，
复制到文末并替换，最后删除全部模板页。

复制页会对 r:embed / r:link / r:id 等关系 ID 做重映射，避免仅 deepcopy XML 时图片/媒体引用指向错误。

依赖模板中每一类版式页的 {占位符} 键集合唯一，且与 JSON 每条 dict 的键集合一致（可通过 key_aliases 映射）。
"""

from __future__ import annotations

import re
from copy import deepcopy
from typing import Any, Dict, FrozenSet, Iterable, List, Mapping, Optional, Set, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import nsuri, qn
from pptx.slide import Slide

from ppt_filler import fill_slide_placeholders, iter_shapes_with_text_frame
from ppt_text_fit import (
    TextFitOptions,
    collect_shape_ids_for_restricted_fit,
    configure_ppt_text_fit_logging,
    fit_slide_text_frames,
)

# 关系命名空间下的属性（r:embed、r:link、r:id 等）值形如 rIdN，需映射到新幻灯片部件的 rels
_REL_NS_PREFIX = "{" + nsuri("r") + "}"

# 示例 JSON 若使用「目录占位符N」，可映射到模板真实键名「目录内容占位符N」
DEFAULT_KEY_ALIASES: Dict[str, str] = {
    "目录占位符1": "目录内容占位符1",
    "目录占位符2": "目录内容占位符2",
    "目录占位符3": "目录内容占位符3",
}


def _collect_r_ids_in_element(root: Any) -> List[str]:
    """收集元素树中所有关系命名空间属性上的 rId* 引用（去重且保持顺序）。"""
    seen: set[str] = set()
    ordered: List[str] = []
    for el in root.iter():
        for attr_name, attr_val in el.attrib.items():
            if not attr_name.startswith(_REL_NS_PREFIX):
                continue
            if not attr_val or not attr_val.startswith("rId"):
                continue
            if attr_val not in seen:
                seen.add(attr_val)
                ordered.append(attr_val)
    return ordered


def _remap_copied_relationships(source_part: Any, dest_part: Any, root: Any) -> None:
    """
    deepcopy 的形状仍携带源页上的 rId，与目标 slide 部件的 _rels 不一致会导致图片/图表丢失。
    按源部件解析每个旧 rId，在目标部件上 relate_to 同一目标，再把 XML 中的 rId 替换为新值。
    """
    old_ids = _collect_r_ids_in_element(root)
    id_map: Dict[str, str] = {}
    for old in old_ids:
        if old not in source_part.rels:
            continue
        rel = source_part.rels[old]
        if rel.is_external:
            new_id = dest_part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_id = dest_part.relate_to(rel.target_part, rel.reltype)
        id_map[old] = new_id
    if not id_map:
        return
    for el in root.iter():
        for attr_name, attr_val in list(el.attrib.items()):
            if attr_val in id_map:
                el.set(attr_name, id_map[attr_val])


def _copy_slide_background(src: Slide, dst: Slide) -> None:
    """将源幻灯片的背景复制到目标幻灯片。bg 元素在 cSld/sld/bg 下，通过迭代查找。"""
    # 1. 找源 bg 元素
    src_bg_elem = None
    for elem in src.element.iter():
        if elem.tag == qn("p:bg"):
            src_bg_elem = elem
            break
    if src_bg_elem is None:
        return  # 无本地背景，继续使用布局/母版继承背景

    # 2. 找目标 spTree（bg 必须插在 spTree 之前）
    dst_spTree = None
    for elem in dst.element.iter():
        if elem.tag == qn("p:spTree"):
            dst_spTree = elem
            break
    if dst_spTree is None:
        return

    # 3. 删除目标已有的 bg 节点（如有）
    for elem in list(dst.element.iter()):
        if elem.tag == qn("p:bg"):
            elem.getparent().remove(elem)

    # 4. 插入复制的 bg 节点
    copied_bg = deepcopy(src_bg_elem)
    dst_spTree.addprevious(copied_bg)


def duplicate_slide(pres: Presentation, slide: Slide) -> Slide:
    """复制一页幻灯片（deepcopy 形状并修复 r:embed 等关系，避免图片/图表引用断裂）。"""
    layout = slide.slide_layout
    new_slide = pres.slides.add_slide(layout)
    for shape in list(new_slide.shapes):
        shape.element.getparent().remove(shape.element)
    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    _remap_copied_relationships(slide.part, new_slide.part, new_slide.element)
    _copy_slide_background(slide, new_slide)
    return new_slide


def delete_slide(pres: Presentation, slide: Slide) -> None:
    """从演示文稿中删除指定页（先移 sldId 再 drop 关系）。"""
    r_id = None
    for rel in pres.part.rels.values():
        if rel.target_part == slide.part:
            r_id = rel.rId
            break
    if not r_id:
        raise ValueError("无法在演示文稿关系中定位该幻灯片")
    sid = slide.slide_id
    sld_id_lst = pres.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        if sld_id.id == sid:
            sld_id_lst.remove(sld_id)
            break
    else:
        raise ValueError("无法在 sldIdLst 中定位该幻灯片")
    pres.part.drop_rel(r_id)


def placeholder_keys_on_slide(slide: Slide) -> FrozenSet[str]:
    """收集该页全部 {key} 中的 key（含组合形状内文本框）。"""
    texts: List[str] = []
    for shape in iter_shapes_with_text_frame(slide.shapes):
        if shape.text:
            texts.append(shape.text)
    full = "".join(texts)
    return frozenset(re.findall(r"\{([^}]+)\}", full))


def _normalize_item(
    item: Mapping[str, Any], key_aliases: Optional[Mapping[str, str]]
) -> Dict[str, Any]:
    """将输入 dict 的键按别名映射为模板中的键名。"""
    aliases = key_aliases or {}
    return {aliases.get(k, k): v for k, v in item.items()}


def _build_template_keymap(
    pres: Presentation, template_count: int
) -> Dict[FrozenSet[str], Slide]:
    """前 template_count 页为模板页：键集合 -> 代表 Slide（后者用于 duplicate）。"""
    keymap: Dict[FrozenSet[str], Slide] = {}
    for i in range(template_count):
        slide = pres.slides[i]
        ks = placeholder_keys_on_slide(slide)
        if not ks:
            continue
        if ks in keymap:
            raise ValueError(
                f"模板第 {i + 1} 页与前面某页的占位符键集合重复: {sorted(ks)}"
            )
        keymap[ks] = slide
    return keymap


def generate_dynamic_ppt(
    template_path: str,
    output_path: str,
    pages: Iterable[Mapping[str, Any]],
    *,
    key_aliases: Optional[Mapping[str, str]] = None,
    text_fit: Optional[TextFitOptions] = None,
) -> Tuple[int, int]:
    """
    按 pages 顺序生成 PPT：每页 dict 的键集合需匹配某一模板页的占位符键集合。

    Args:
        template_path: 含「模板页」的 pptx（通常全部版式页连续排在文件前部）
        output_path: 输出路径
        pages: 列表/可迭代，每项为占位符 -> 值
        key_aliases: 输入键名 -> 模板键名，会与 DEFAULT_KEY_ALIASES 合并（传入优先覆盖）
        text_fit: 非 None 且 enabled=True 时，在占位符替换后对每页执行按框字号适配（见 ppt_text_fit）。
            若 text_fit.only_fit_page_title_content_blocks=True，则仅在 fill 前根据模板占位符
            收集「内容占位符」所在文本框，避免标题等被缩小。

    Returns:
        (生成的内容页数量, 累计替换占位符次数)
    """
    merged_aliases: Dict[str, str] = {**DEFAULT_KEY_ALIASES}
    if key_aliases:
        merged_aliases.update(dict(key_aliases))

    prs = Presentation(template_path)
    template_count = len(prs.slides)
    if template_count == 0:
        raise ValueError("模板无幻灯片")

    keymap = _build_template_keymap(prs, template_count)
    total_replacements = 0
    n_out = 0

    for raw in pages:
        item = _normalize_item(raw, merged_aliases)
        ks = frozenset(item.keys())
        if ks not in keymap:
            available = [sorted(k) for k in keymap.keys()]
            raise KeyError(
                f"无匹配模板页，键集合为 {sorted(ks)}；\n"
                f"可用模板键集合 ({len(available)}个): {available}"
            )
        src = keymap[ks]
        new_slide = duplicate_slide(prs, src)
        shape_ids_to_fit: Optional[Set[int]] = None
        if text_fit is not None and text_fit.only_fit_page_title_content_blocks:
            shape_ids_to_fit = collect_shape_ids_for_restricted_fit(
                new_slide, item, text_fit
            )
        total_replacements += fill_slide_placeholders(new_slide, item)
        if text_fit is not None:
            fit_slide_text_frames(new_slide, text_fit, shape_ids_to_fit=shape_ids_to_fit)
        n_out += 1

    for _ in range(template_count):
        delete_slide(prs, prs.slides[0])

    prs.save(output_path)
    return n_out, total_replacements


if __name__ == "__main__":
    # dw5.pptx 模板键集合（实际检测）:
    # 1: ['主标题']
    # 2: ['目录内容占位符1', '目录内容占位符2', '目录内容占位符3']
    # 3: ['章节序号', '章节标题']
    # 4: ['内容占位符1', '小标题占位符1', '页面标题（1框）']
    # 5: ['内容占位符1', '内容占位符2', '小标题占位符1', '小标题占位符2', '页面标题（2框）']
    # 6: ['内容占位符1', '内容占位符2', '内容占位符3', '小标题占位符1', '小标题占位符2', '小标题占位符3']  # 3框无页面标题
    # 7: ['内容占位符1', '内容占位符2', '内容占位符3', '内容占位符4', '小标题占位符1', '小标题占位符2', '小标题占位符3', '小标题占位符4', '页面标题（4框）']
    # 8: ['结束语']
    configure_ppt_text_fit_logging()
    generate_dynamic_ppt(
        template_path="templates/会议模板1.pptx",
        output_path="output.pptx",
        text_fit=TextFitOptions(only_fit_page_title_content_blocks=True),
        pages=[],
    )

