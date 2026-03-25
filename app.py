#!/usr/bin/env python3
"""
Flask 应用：接收表单/JSON（含 pages 列表），调用 ppt_dynamic_filler 生成 PPT，
将文件保存到本地 output 目录，提供前端页面和下载/预览功能。
支持通过 Qwen API 从文档生成 PPT 内容。
"""


from __future__ import annotations

from dotenv import load_dotenv
load_dotenv()
import json
import os
import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from flask import Flask, jsonify, request, send_file, render_template, Response
from ppt_dynamic_filler import generate_dynamic_ppt
from pptx import Presentation


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
TEMPLATE_DIR = BASE_DIR / "templates"
OUTPUT_DIR.mkdir(exist_ok=True)

# 默认模板路径
_DEFAULT_REL = os.environ.get("PPT_DEFAULT_TEMPLATE", "templates/会议模板2.pptx")
_ALT_REL = "templates/会议模板1.pptx"

# 默认提示词
_DEFAULT_PROMPT = """## 角色
你是一个善于根据文字素材总结出格式化数据的摘要专家。

## 要求与限制
1.请总结用户发给你的素材内容，输出列表嵌套字典的json数据;
2.先确定标题和三个章节标题（固定生成三个章节），目录内容和章节标题是相对应的，三个目录内容占位符对应了三个章节标题；
3.在每个对应的章节中生成内容页的字典，按内容分的点数确定内容页的字典格式。
包含页面标题（1框）、页面标题（2框）等的这些键是内容页，无需严格按照下面的数量生成字典，但内容页的字典数量要与内容分的点数一致。
例如，某章节的一个要点只分一点就用{"页面标题（1框）": "", "小标题占位符1": "", "内容占位符1": ""},
分两点就用{"页面标题（2框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符2": ""},
分三点就用{"页面标题（3框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符3": ""}
分四点就用{"页面标题（4框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符3": "", "小标题占位符4": "", "内容占位符4": ""}
4.最终输出结果的格式大致如下，务必要生成完整的JSON字符串：
[
    {"主标题": ""},
    {"目录内容占位符1": "", "目录内容占位符2": "", "目录内容占位符3": ""},
    {"章节序号": "01", "章节标题": ""},
    {"页面标题（1框）": "", "小标题占位符1": "", "内容占位符1": ""},
    {"页面标题（2框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符2": ""},
    {"章节序号": "02", "章节标题": ""},
    {"页面标题（2框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符2": ""},
    {"章节序号": "03", "章节标题": ""},
    {"页面标题（3框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符3": ""},
    {"页面标题（4框）": "", "小标题占位符1": "", "内容占位符1": "", "小标题占位符2": "", "内容占位符3": "", "小标题占位符4": "", "内容占位符4": ""},
    {"结束语": "谢谢！"}
]"""

BASE_DIR = Path(__file__).resolve().parent
OUTPUT_DIR = BASE_DIR / "output"
TEMPLATE_DIR = BASE_DIR / "templates"
OUTPUT_DIR.mkdir(exist_ok=True)

# 默认模板路径
_DEFAULT_REL = os.environ.get("PPT_DEFAULT_TEMPLATE", "templates/会议模板2.pptx")
_ALT_REL = "templates/会议模板1.pptx"


def _resolve_template_path(explicit: Optional[str]) -> Path:
    """解析模板文件路径（支持相对项目根目录）。"""
    if explicit:
        p = Path(explicit)
        return p if p.is_absolute() else (BASE_DIR / p)
    primary = BASE_DIR / _DEFAULT_REL
    if primary.exists():
        return primary
    alt = BASE_DIR / _ALT_REL
    if alt.exists():
        return alt
    return primary


def _build_output_filename(prefix: Optional[str] = None) -> str:
    """生成输出文件名，格式：{前缀}_{日期}_{随机ID}.pptx"""
    day = datetime.now().strftime("%Y%m%d")
    base = (prefix or "ppt-generated").strip("/").strip() or "ppt-generated"
    return f"{base}_{day}_{uuid.uuid4().hex}.pptx"


def _get_available_templates() -> List[Dict[str, str]]:
    """获取可用的模板文件列表。"""
    templates = []
    for p in sorted(TEMPLATE_DIR.glob("*.pptx")):
        templates.append({
            "name": p.stem,
            "path": f"templates/{p.name}",
            "display_name": p.stem
        })
    return templates


def _get_generated_files() -> List[Dict[str, Any]]:
    """获取已生成的 PPT 文件列表。"""
    files = []
    for p in sorted(OUTPUT_DIR.glob("*.pptx"), key=lambda x: x.stat().st_mtime, reverse=True):
        files.append({
            "name": p.name,
            "size": p.stat().st_size,
            "created": datetime.fromtimestamp(p.stat().st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
            "modified": datetime.fromtimestamp(p.stat().st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
        })
    return files


def _validate_pages(data: Any) -> Tuple[bool, str, List[Dict[str, Any]]]:
    """校验 data 是否为非空 dict 列表。"""
    if not isinstance(data, list):
        return False, "data 必须为 JSON 数组", []
    if len(data) == 0:
        return False, "data 不能为空列表", []
    out: List[Dict[str, Any]] = []
    for i, item in enumerate(data):
        if not isinstance(item, dict):
            return False, f"data[{i}] 必须为对象", []
        out.append(item)
    return True, "", out


def _extract_ppt_content(file_path: Path) -> List[Dict[str, Any]]:
    """提取 PPT 的每页内容用于预览。"""
    from ppt_filler import iter_shapes_with_text_frame

    prs = Presentation(str(file_path))
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in iter_shapes_with_text_frame(slide.shapes):
            if shape.text and shape.text.strip():
                texts.append(shape.text.strip())

        slides_data.append({
            "index": slide_idx,
            "texts": texts
        })

    return slides_data


def _read_text_file(file_path: Path) -> str:
    """读取文本文件内容。"""
    encodings = ['utf-8', 'gbk', 'gb18030']
    for encoding in encodings:
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _extract_doc_text(doc_path: str) -> Optional[str]:
    """用 pypandoc 提取 .doc/.docx 文本，失败返回 None。"""
    import pypandoc
    try:
        text = pypandoc.convert_file(doc_path, 'plain', format='doc')
        return text if text else None
    except Exception:
        return None


def _call_qwen_api(file_content: str, prompt: str, model: str) -> Dict[str, Any]:
    """调用 Qwen API 生成 PPT JSON 数据。"""
    try:
        from openai import OpenAI

        api_key = os.getenv("DASHSCOPE_API_KEY")
        if not api_key:
            return {"ok": False, "message": "未配置 DASHSCOPE_API_KEY 环境变量"}

        client = OpenAI(
            api_key=api_key,
            base_url="https://dashscope.aliyuncs.com/compatible-mode/v1",
        )

        messages = [
            {"role": "system", "content": prompt},
            {"role": "user", "content": file_content}
        ]

        completion = client.chat.completions.create(
            model=model,
            messages=messages,
        )

        content = completion.model_dump_json()
        content_data = json.loads(content).get("choices", [{}])[0].get("message", {}).get("content", "")

        # 清理 JSON 格式（code fence 是在 content_data 即 AI 返回的文字中）
        raw = content_data
        if raw.startswith("```json\n"):
            content_data = raw[8:]
        elif raw.startswith("```json"):
            content_data = raw[7:]
        if content_data.endswith("```\n"):
            content_data = content_data[:-4]
        elif content_data.endswith("```"):
            content_data = content_data[:-3]

        try:
            result = json.loads(content_data)
            return {"ok": True, "data": result}
        except json.JSONDecodeError:
            return {"ok": False, "message": "API 返回的 JSON 格式错误"}

    except Exception as e:
        return {"ok": False, "message": f"调用 Qwen API 失败: {str(e)}"}


def create_app() -> Flask:
    """创建并配置 Flask 应用。"""
    app = Flask(__name__, template_folder='templates')

    @app.get("/")
    def index() -> Any:
        """前端首页。"""
        templates = _get_available_templates()
        files = _get_generated_files()
        return render_template("index.html", templates=templates, files=files)

    @app.get("/health")
    def health() -> Any:
        """健康检查。"""
        return jsonify({"status": "ok"})

    @app.get("/api/templates")
    def get_templates() -> Any:
        """获取可用模板列表。"""
        templates = _get_available_templates()
        return jsonify({"ok": True, "templates": templates})

    @app.get("/api/prompt/default")
    def get_default_prompt() -> Any:
        """获取默认提示词。"""
        return jsonify({"ok": True, "prompt": _DEFAULT_PROMPT})

    @app.post("/api/ai/generate")
    def ai_generate() -> Any:
        """通过 AI 生成 PPT JSON 数据。"""
        body = request.get_json(silent=True)
        if not isinstance(body, dict):
            return jsonify({"ok": False, "message": "请求体须为 JSON 对象"}), 400

        file_content = body.get("content", "")
        prompt = body.get("prompt", _DEFAULT_PROMPT)
        model = body.get("model", "qwen-long")

        if not file_content:
            return jsonify({"ok": False, "message": "请提供文档内容"}), 400

        # 调用 Qwen API
        result = _call_qwen_api(file_content, prompt, model)
        return jsonify(result)

    @app.get("/api/files")
    def get_files() -> Any:
        """获取已生成的 PPT 文件列表。"""
        files = _get_generated_files()
        return jsonify({"ok": True, "files": files})

    @app.post("/api/ppt/generate")
    def generate_ppt() -> Any:
        """
        请求 JSON:
          - data: 必填，与 ppt_dynamic_filler.generate_dynamic_ppt 的 pages 一致
          - template: 可选，模板 pptx 相对或绝对路径
          - filename_prefix: 可选，输出文件名前缀（默认 ppt-generated）
          - key_aliases: 可选，占位符键别名映射对象

        响应 JSON: filename, download_url, preview_url, slides, replacements；错误时 message + HTTP 4xx/5xx。
        """
        body = request.get_json(silent=True)
        if not isinstance(body, dict):
            return jsonify({"ok": False, "message": "请求体须为 JSON 对象"}), 400
        if "data" not in body:
            return jsonify({"ok": False, "message": "缺少字段 data"}), 400

        ok, err, pages = _validate_pages(body["data"])
        if not ok:
            return jsonify({"ok": False, "message": err}), 400

        template_path = _resolve_template_path(
            body.get("template") if isinstance(body.get("template"), str) else None
        )
        if not template_path.is_file():
            return (
                jsonify(
                    {
                        "ok": False,
                        "message": f"模板文件不存在: {template_path}",
                    }
                ),
                400,
            )

        filename_prefix = body.get("filename_prefix")
        if filename_prefix is not None and not isinstance(filename_prefix, str):
            return jsonify({"ok": False, "message": "filename_prefix 须为字符串"}), 400

        key_aliases = body.get("key_aliases")
        if key_aliases is not None and not isinstance(key_aliases, dict):
            return jsonify({"ok": False, "message": "key_aliases 须为对象"}), 400

        try:
            output_filename = _build_output_filename(filename_prefix)
            output_path = OUTPUT_DIR / output_filename

            slides, replacements = generate_dynamic_ppt(
                str(template_path),
                str(output_path),
                pages,
                key_aliases=key_aliases,
            )

            return jsonify(
                {
                    "ok": True,
                    "filename": output_filename,
                    "download_url": f"/api/files/{output_filename}/download",
                    "preview_url": f"/api/files/{output_filename}/preview",
                    "slides": slides,
                    "replacements": replacements,
                }
            )
        except KeyError as e:
            return jsonify({"ok": False, "message": f"占位符与模板不匹配: {e!s}"}), 400
        except (ValueError, OSError) as e:
            return jsonify({"ok": False, "message": str(e)}), 400
        except Exception as e:  # noqa: BLE001 — 生成未知错误统一返回 500
            return jsonify({"ok": False, "message": f"服务异常: {e!s}"}), 500

    @app.get("/api/files/<filename>/download")
    def download_file(filename: str) -> Any:
        """下载 PPT 文件。"""
        file_path = OUTPUT_DIR / filename
        if not file_path.is_file():
            return jsonify({"ok": False, "message": "文件不存在"}), 404
        return send_file(str(file_path), as_attachment=True, download_name=filename)

    @app.get("/api/files/<filename>/preview")
    def preview_file(filename: str) -> Any:
        """预览 PPT 文件（返回 HTML 预览页面）。"""
        file_path = OUTPUT_DIR / filename
        if not file_path.is_file():
            return jsonify({"ok": False, "message": "文件不存在"}), 404

        slides_data = _extract_ppt_content(file_path)
        return render_template("preview.html", filename=filename, slides=slides_data)

    @app.post("/api/upload/doc")
    def upload_doc() -> Any:
        """上传 .doc 文件，后端提取文本内容返回。"""
        if 'file' not in request.files:
            return jsonify({"ok": False, "message": "未找到上传文件"}), 400

        file = request.files['file']
        if not file.filename:
            return jsonify({"ok": False, "message": "文件名为空"}), 400

        ext = file.filename.lower().split('.')[-1]
        if ext != 'doc':
            return jsonify({"ok": False, "message": "仅支持 .doc 文件"}), 400

        # 保存到临时文件
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name

        try:
            text = _extract_doc_text(tmp_path)
        finally:
            Path(tmp_path).unlink(missing_ok=True)

        if text is None:
            return jsonify({"ok": False, "message": ".doc 文件解析失败，请确保已安装 pandoc: https://pandoc.org/installing.html"}), 500

        return jsonify({"ok": True, "content": text})

    @app.post("/api/files/<filename>/delete")
    def delete_file(filename: str) -> Any:
        """删除 PPT 文件。"""
        file_path = OUTPUT_DIR / filename
        if not file_path.is_file():
            return jsonify({"ok": False, "message": "文件不存在"}), 404
        try:
            file_path.unlink()
            return jsonify({"ok": True, "message": "文件已删除"})
        except OSError as e:
            return jsonify({"ok": False, "message": f"删除失败: {e!s}"}), 500

    return app


if __name__ == "__main__":
    app = create_app()
    port = int(os.environ.get("PORT", "5051"))
    app.run(host="0.0.0.0", port=port, debug=os.environ.get("FLASK_DEBUG") == "1")
