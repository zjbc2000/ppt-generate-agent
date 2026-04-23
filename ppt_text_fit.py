#!/usr/bin/env python3
"""
占位符填充后：按文本框几何尺寸缩小字号（python-pptx TextFrame.fit_text），
可选最小字号与字库路径。整框会统一为同一字体/字号（见 fit_text 行为）。
"""

from __future__ import annotations

import logging
import os
import sys
from dataclasses import dataclass, field
from typing import FrozenSet, List, Literal, Mapping, Optional, Set

from pptx.text.fonts import FontFiles
from pptx.text.layout import _BinarySearchTree, _rendered_size
from pptx.util import Pt

from ppt_filler import iter_shapes_with_text_frame

logger = logging.getLogger(__name__)

OnUnderflowMin = Literal["warn", "none"]


def configure_ppt_text_fit_logging() -> None:
    """
    为 ppt_text_fit logger 挂载 StreamHandler，使 INFO 在控制台可见。
    根 logger 常为 WARNING 时，默认看不到「适配后字号」类 INFO；调用一次即可（幂等）。
    """
    if logger.handlers:
        return
    logger.setLevel(logging.INFO)
    h = logging.StreamHandler()
    h.setLevel(logging.INFO)
    h.setFormatter(logging.Formatter("%(levelname)s [ppt_text_fit] %(message)s"))
    logger.addHandler(h)
    logger.propagate = False

# 仅对「含章节页面标题」的版式页上的「正文内容」占位框做字号适配（可经 TextFitOptions 覆盖）
DEFAULT_PAGE_TITLE_TRIGGER_KEYS: FrozenSet[str] = frozenset(
    ("页面标题（1框）", "页面标题（2框）", "页面标题（3框）", "页面标题（4框）")
)
DEFAULT_CONTENT_FIT_KEYS: FrozenSet[str] = frozenset(
    ("内容占位符1", "内容占位符2", "内容占位符3", "内容占位符4")
)


@dataclass
class TextFitOptions:
    """控制 fit_slide_text_frames 行为。"""

    enabled: bool = True
    """为 False 时不做任何适配。"""
    font_file: Optional[str] = None
    """若指定，优先用于 PIL 测宽（中文建议显式传入 .ttf/.ttc 路径）。"""
    font_file_map: Optional[Mapping[str, str]] = None
    """字体显示名（如「微软雅黑」）到字库文件路径的映射。
    键 \"*\" 或 \"__default__\" 表示任意未单独映射的模板字体共用该字库文件。"""
    fallback_font_file: Optional[str] = None
    """当模板字体名在系统字库索引中不存在（FontFiles.find 失败）时，用于 PIL 测宽的 .ttf/.otf 路径。
    与 font_file_map 的通配键二选一或同用；建议指向与版式接近的中文无衬线字库。"""
    min_pt: float = 8.0
    """fit 后若字号低于此值，则抬升到该值（可能溢出框，见 on_underflow_min）。"""
    max_pt_cap: Optional[float] = None
    """若设置，模板字号上限不超过该值（pt）。"""
    on_underflow_min: OnUnderflowMin = "warn"
    """当因 min_pt 抬升而可能溢出时是否打日志。"""
    only_fit_page_title_content_blocks: bool = False
    """若为 True：仅当 item 的键与 page_title_trigger_keys 有交集时，本页才做 fit；
    且只处理模板文本中仍含 {内容占位符N}（见 content_fit_keys）且该键在 item 中的那些文本框。
    须在 fill 之前根据模板占位符收集目标 shape（由 generate_dynamic_ppt 传入 shape_ids_to_fit）。"""
    page_title_trigger_keys: FrozenSet[str] = field(
        default_factory=lambda: frozenset(DEFAULT_PAGE_TITLE_TRIGGER_KEYS)
    )
    content_fit_keys: FrozenSet[str] = field(
        default_factory=lambda: frozenset(DEFAULT_CONTENT_FIT_KEYS)
    )


def _shape_label(shape: object) -> str:
    try:
        name = getattr(shape, "name", None)
        return str(name) if name else "unnamed"
    except Exception:
        return "unnamed"


def _first_paragraph_first_run(shape: object):
    tf = shape.text_frame
    for para in tf.paragraphs:
        if para.runs:
            return para.runs[0]
    return None


def _log_fit_result(shape: object, label: str, template_max_pt: int) -> None:
    """适配流程结束后打 INFO，便于确认是否缩小字号（未配置 logging 时可能仍不显示）。"""
    run = _first_paragraph_first_run(shape)
    if run is None or run.font.size is None:
        return
    fitted = float(run.font.size.pt)
    logger.info(
        "文本框「%s」字号: 模板上限 %dpt → 适配后 %.1fpt",
        label,
        template_max_pt,
        fitted,
    )


def _apply_min_pt_to_frame(shape: object, min_pt: float) -> None:
    """将文本框内各 run 字号设为 min_pt（不处理 endParaRPr 等边角）。"""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(min_pt)


def _is_existing_file(path: Optional[str]) -> bool:
    return bool(path and os.path.isfile(path))


def _raw_text_frame_extents_emu(shape: object) -> tuple[int, int]:
    """与 TextFrame._extents 相同公式；边距大于形状时可能为负（模板常见问题）。"""
    tf = shape.text_frame
    parent = tf._parent
    cw = int(parent.width)
    ch = int(parent.height)
    ml, mr = int(tf.margin_left), int(tf.margin_right)
    mt, mb = int(tf.margin_top), int(tf.margin_bottom)
    return cw - ml - mr, ch - mt - mb


def _clamped_text_frame_extents_emu(shape: object) -> tuple[int, int]:
    """
    用于测宽/折行的有效宽高（EMU），保证为正。

    当边距之和大于 cx/cy 时，PowerPoint 仍会显示文本，但按公式会得到负值；
    此处用形状尺寸的近似比例作为测宽区域，避免跳过适配或库内部 TextFitter 异常。
    """
    tf = shape.text_frame
    parent = tf._parent
    cw = int(parent.width)
    ch = int(parent.height)
    ml, mr = int(tf.margin_left), int(tf.margin_right)
    mt, mb = int(tf.margin_top), int(tf.margin_bottom)
    w = cw - ml - mr
    h = ch - mt - mb
    if cw <= 0:
        w = 457200  # 0.5 英寸，兜底
    elif w <= 0:
        w = max(1, int(cw * 0.88))
    if ch <= 0:
        h = 457200
    elif h <= 0:
        h = max(1, int(ch * 0.88))
    return max(1, w), max(1, h)


def _wrap_paragraph_chars(
    paragraph: str, width_emu: int, pt: int, font_file: str
) -> Optional[List[str]]:
    """
    按字符折行（弥补 python-pptx TextFitter 仅按空白分词，中文无空格时 _break_line 返回 None 的问题）。
    """
    if not paragraph:
        return []
    lines: List[str] = []
    current = ""
    for ch in paragraph:
        trial = current + ch
        tw = _rendered_size(trial, pt, font_file)[0]
        if tw <= width_emu:
            current = trial
            continue
        if current:
            lines.append(current)
            current = ch
            cw = _rendered_size(current, pt, font_file)[0]
            if cw > width_emu:
                return None
        else:
            return None
    if current:
        lines.append(current)
    return lines


def _text_fits_charwrap(
    text: str, width_emu: int, height_emu: int, pt: int, font_file: str
) -> bool:
    """整段文本（含 \\n 分段）在字号 pt 下是否不超出宽、高。"""
    if width_emu <= 0 or height_emu <= 0:
        return False
    all_lines: List[str] = []
    for part in text.split("\n"):
        sub = _wrap_paragraph_chars(part, width_emu, pt, font_file)
        if sub is None:
            return False
        all_lines.extend(sub)
    line_h = _rendered_size("Ty", pt, font_file)[1]
    return line_h * len(all_lines) <= height_emu


def _best_fit_font_size_charwrap(
    text: str,
    width_emu: int,
    height_emu: int,
    max_size: int,
    font_file: str,
) -> int:
    """与 TextFitter 类似，但用逐字折行判断是否能放入。"""
    sizes = _BinarySearchTree.from_ordered_sequence(range(1, int(max_size) + 1))

    def predicate(pt: int) -> bool:
        return _text_fits_charwrap(text, width_emu, height_emu, pt, font_file)

    best = sizes.find_max(predicate)
    return int(best) if best is not None else 1


def _apply_fit_text_charwrap_fallback(
    shape: object,
    font_family: str,
    max_size: int,
    bold: bool,
    italic: bool,
    font_file: str,
    options: TextFitOptions,
    *,
    width_emu: Optional[int] = None,
    height_emu: Optional[int] = None,
) -> List[str]:
    """
    当 TextFrame.fit_text / TextFitter 内部 unpack None，或边距溢出导致库无法安全 fit 时，
    用逐字折行重新计算字号并 _apply_fit。
    """
    local_warns: List[str] = []
    tf = shape.text_frame
    text = tf.text
    if width_emu is not None and height_emu is not None:
        w, h = width_emu, height_emu
    else:
        w, h = _clamped_text_frame_extents_emu(shape)
    # 防护：极小有效区域直接跳过，避免算出荒谬的1pt
    min_emu = 3600  # ≈1mm，低于此值认为模板设计异常，不适配
    if w < min_emu or h < min_emu:
        local_warns.append(
            f"文本框「{_shape_label(shape)}」有效区域过小（{w}x{h} EMU），跳过适配。"
        )
        return local_warns

    best = _best_fit_font_size_charwrap(text, w, h, max_size, font_file)
    best = max(1, min(best, max_size))
    fitted_pt = float(best)

    tf._apply_fit(font_family, best, bool(bold), bool(italic))

    if fitted_pt < options.min_pt:
        _apply_min_pt_to_frame(shape, options.min_pt)
        if options.on_underflow_min == "warn":
            msg = (
                f"文本框「{_shape_label(shape)}」逐字折行适配后字号 {fitted_pt:.1f}pt 低于 "
                f"min_pt={options.min_pt}，已抬升到 {options.min_pt}pt，可能溢出框外。"
            )
            logger.warning(msg)
            local_warns.append(msg)

    return local_warns


def _resolve_font_file_for_fit(
    font_family: str,
    bold: bool,
    italic: bool,
    options: TextFitOptions,
) -> Optional[str]:
    """
    解析用于 TextFitter / PIL 的字库路径。

    python-pptx 的 FontFiles 仅索引部分 .otf/.ttf；模板里「汉仪」「方正」等名常无法命中，
    此时必须用 font_file_map、通配键或 fallback_font_file，否则会 KeyError。
    """
    if _is_existing_file(options.font_file):
        return options.font_file

    if options.font_file_map:
        m = options.font_file_map
        for key in (font_family, "*", "__default__"):
            if key is None:
                continue
            p = m.get(key)
            if _is_existing_file(p):
                return p

    try:
        return FontFiles.find(font_family, bold, italic)
    except KeyError:
        pass

    if _is_existing_file(options.fallback_font_file):
        return options.fallback_font_file

    for name in (
        "Arial",
        "Calibri",
        "Microsoft YaHei",
        "PingFang SC",
        "Heiti SC",
        "STHeiti",
        "SimHei",
        "SimSun",
    ):
        try:
            return FontFiles.find(name, False, False)
        except KeyError:
            continue

    # FontFiles 只索引 .otf/.ttf，macOS 上大量字体为 .ttc 不会入库，尝试常见绝对路径
    if sys.platform == "darwin":
        darwin_paths = (
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
            "/System/Library/Fonts/Supplemental/Times New Roman.ttf",
            "/Library/Fonts/Arial.ttf",
        )
        for p in darwin_paths:
            if _is_existing_file(p):
                return p

    return None


def collect_shape_ids_for_restricted_fit(
    slide: object,
    item: Mapping[str, Any],
    options: TextFitOptions,
) -> Optional[Set[int]]:
    """
    在 fill_slide_placeholders **之前**调用：根据模板中仍存在的 {{占位符}} 判定哪些 shape 需要 fit。

    Returns:
        None — 未启用 only_fit_page_title_content_blocks，应对全页所有文本框 fit。
        set() — 本页不应 fit 任何框（无触发键或无可匹配占位符）。
        非空 set — 仅对这些 id(shape) 做 fit。
    """
    if not options.only_fit_page_title_content_blocks:
        return None
    item_keys = frozenset(item.keys())
    if not (item_keys & options.page_title_trigger_keys):
        return None
    out: Set[int] = set()
    for shape in iter_shapes_with_text_frame(slide.shapes):
        full = shape.text_frame.text
        for k in options.content_fit_keys:
            if k not in item:
                continue
            if f"{{{k}}}" in full:
                out.add(id(shape))
                break
    if not out:
        return None
    return out


def _fit_single_shape(shape: object, options: TextFitOptions) -> List[str]:
    """对单个带 text_frame 的 shape 做字号适配，返回该形状产生的告警。"""
    warnings: List[str] = []
    raw = shape.text_frame.text
    if not raw or not str(raw).strip():
        return warnings

    run0 = _first_paragraph_first_run(shape)
    if run0 is None:
        return warnings

    font_family = run0.font.name or "Calibri"
    sz = run0.font.size
    base_pt = float(sz.pt) if sz is not None else 18.0
    max_size = int(round(base_pt))
    if options.max_pt_cap is not None:
        max_size = min(max_size, int(round(options.max_pt_cap)))
    max_size = max(1, max_size)

    bold = run0.font.bold if run0.font.bold is not None else False
    italic = run0.font.italic if run0.font.italic is not None else False

    font_file = _resolve_font_file_for_fit(
        font_family, bool(bold), bool(italic), options
    )

    label = _shape_label(shape)
    if font_file is None:
        msg = (
            f"文本框「{label}」无法解析字库文件（模板字体「{font_family}」），"
            f"请设置 TextFitOptions.font_file、font_file_map（含「*」）或 fallback_font_file。"
        )
        logger.warning(msg)
        warnings.append(msg)
        return warnings

    raw_w, raw_h = _raw_text_frame_extents_emu(shape)
    eff_w, eff_h = _clamped_text_frame_extents_emu(shape)

    if raw_w <= 0 or raw_h <= 0:
        logger.debug(
            "文本框「%s」边距大于形状尺寸 (raw=%sx%s EMU)，仅用逐字折行 + 近似区域 %sx%s",
            label,
            raw_w,
            raw_h,
            eff_w,
            eff_h,
        )
        warnings.extend(
            _apply_fit_text_charwrap_fallback(
                shape,
                font_family,
                max_size,
                bold,
                italic,
                font_file,
                options,
                width_emu=eff_w,
                height_emu=eff_h,
            )
        )
        _log_fit_result(shape, label, max_size)
        return warnings

    try:
        shape.text_frame.fit_text(
            font_family=font_family,
            max_size=max_size,
            bold=bool(bold),
            italic=bool(italic),
            font_file=font_file,
        )
    except Exception as e:
        err = str(e)
        if "NoneType" in err and "unpack" in err:
            try:
                warnings.extend(
                    _apply_fit_text_charwrap_fallback(
                        shape,
                        font_family,
                        max_size,
                        bold,
                        italic,
                        font_file,
                        options,
                    )
                )
                return warnings
            except Exception as e2:
                msg = f"文本框「{label}」fit_text 与逐字折行回退均失败: {e2}"
                logger.warning(msg)
                warnings.append(msg)
        else:
            msg = f"文本框「{label}」fit_text 失败: {e}"
            logger.warning(msg)
            warnings.append(msg)
        return warnings

    run_after = _first_paragraph_first_run(shape)
    if run_after is None or run_after.font.size is None:
        return warnings
    fitted_pt = float(run_after.font.size.pt)

    if fitted_pt < options.min_pt:
        _apply_min_pt_to_frame(shape, options.min_pt)
        if options.on_underflow_min == "warn":
            msg = (
                f"文本框「{label}」适配后字号 {fitted_pt:.1f}pt 低于 min_pt={options.min_pt}，"
                f"已抬升到 {options.min_pt}pt，可能溢出框外。"
            )
            logger.warning(msg)
            warnings.append(msg)

    _log_fit_result(shape, label, max_size)

    return warnings


def fit_slide_text_frames(
    slide: object,
    options: TextFitOptions,
    *,
    shape_ids_to_fit: Optional[Set[int]] = None,
) -> List[str]:
    """
    对幻灯片内非空文本框调用 TextFrame.fit_text（或逐字折行回退），并按需施加 min_pt 下限。

    shape_ids_to_fit:
        None — 遍历所有带文本的形状（未启用按占位符筛选时由调用方传 None）。
        set() — 不处理任何形状（例如本页无「页面标题」触发键）。
        非空 set — 仅处理 id(shape) 落在集合内的形状（须在 fill 前收集）。

    fit_text 会将整框统一为同一 font_family / 字号 / bold / italic；
    中文测宽建议配置 font_file 或 font_file_map。

    Returns:
        人类可读告警文案列表（便于调用方收集或展示）。
    """
    warnings: List[str] = []
    if not options.enabled:
        return warnings
    if shape_ids_to_fit is not None and len(shape_ids_to_fit) == 0:
        logger.info("本页未做字号适配（无匹配文本框或非页面标题版式）。")
        return warnings

    processed_ids: set[int] = set()
    for shape in iter_shapes_with_text_frame(slide.shapes):
        shape_id = id(shape)
        if shape_id in processed_ids:
            continue
        if shape_ids_to_fit is not None and shape_id not in shape_ids_to_fit:
            continue
        processed_ids.add(shape_id)
        warnings.extend(_fit_single_shape(shape, options))

    return warnings
