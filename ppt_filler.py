#!/usr/bin/env python3
"""
PPT模板填充引擎 - 基于占位符替换

简单的占位符填充逻辑，支持 {key} 格式的占位符；遍历组合形状内文本框。
"""

from copy import deepcopy
from typing import Any, Dict, Iterable

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape


def iter_shapes_with_text_frame(
    shapes: Iterable[BaseShape],
) -> Iterable[BaseShape]:
    """递归遍历形状树，产出带 text_frame 的形状（含组合内子形状）。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_text_frame(shape.shapes)  # type: ignore[attr-defined]
        elif shape.has_text_frame:
            yield shape


def fill_slide_placeholders(slide: Any, data: Dict[str, Any]) -> int:
    """替换单页内所有 {key} 占位符，返回替换次数。"""
    replaced_count = 0
    for shape in iter_shapes_with_text_frame(slide.shapes):
        for paragraph in shape.text_frame.paragraphs:
            full_text = paragraph.text
            new_text = full_text
            has_replacement = False
            for key, value in data.items():
                placeholder = f"{{{key}}}"
                if placeholder in new_text:
                    new_text = new_text.replace(placeholder, str(value))
                    has_replacement = True
                    replaced_count += 1

            if not has_replacement:
                continue

            if paragraph.runs:
                first_run = paragraph.runs[0]
                saved_font = deepcopy(first_run.font)
                for run in paragraph.runs:
                    run.text = ""
                paragraph.runs[0].text = new_text
                run = paragraph.runs[0]
                if saved_font.name:
                    run.font.name = saved_font.name
                if saved_font.size:
                    run.font.size = saved_font.size
                if saved_font.bold is not None:
                    run.font.bold = saved_font.bold
                if saved_font.italic is not None:
                    run.font.italic = saved_font.italic
                try:
                    if saved_font.color and saved_font.color.rgb:
                        run.font.color.rgb = saved_font.color.rgb
                except AttributeError:
                    pass
            else:
                paragraph.add_run().text = new_text

    return replaced_count


def fill_ppt_template(template_path: str, output_path: str, data: Dict[str, Any]) -> int:
    """
    填充PPT模板中的占位符

    占位符格式: {key}

    Args:
        template_path: 模板文件路径
        output_path: 输出文件路径
        data: 替换数据字典

    Returns:
        替换的占位符数量
    """
    prs = Presentation(template_path)
    replaced_count = 0
    for slide in prs.slides:
        replaced_count += fill_slide_placeholders(slide, data)
    prs.save(output_path)
    return replaced_count
